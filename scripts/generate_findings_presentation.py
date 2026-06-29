#!/usr/bin/env python3
"""ISAAI Consolidated Findings Presentation Generator.

Reads all processed test runs from tests/runs/ and generates a single
PowerPoint presentation consolidating all validation findings.
Replaces the old per-run presentation approach with one unified deck.
"""

import sys
import re
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx is required. Install it with: pip install python-pptx")
    sys.exit(1)

# Setup paths relative to the script location
ROOT = Path(__file__).resolve().parents[1]
TESTS_DIR = ROOT / "tests"
RUNS_DIR = TESTS_DIR / "runs"
PRESENTATIONS_DIR = ROOT / "presentations"

THRESHOLD_LIMIT = 50

# Style guide colors matching the project charter
PRIMARY = RGBColor(0x00, 0x33, 0x5B)      # Slate Blue
SECONDARY = RGBColor(0x00, 0x50, 0x8F)    # Active Blue
MUTED = RGBColor(0x66, 0x66, 0x66)        # Neutral Gray
ALERT_RED = RGBColor(0xD9, 0x38, 0x3A)    # Warning Red
PASS_GREEN = RGBColor(0x2E, 0x7D, 0x32)   # Pass Green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def parse_runs():
    """Scans tests/runs/ directories, parsing run summaries."""
    runs = []
    if not RUNS_DIR.exists():
        print(f"Error: {RUNS_DIR} does not exist. Run scripts/simulate_runs.py first.")
        sys.exit(1)

    for run_path in sorted(RUNS_DIR.glob("run_*")):
        summary_path = run_path / "summary.txt"
        if not summary_path.exists():
            continue

        run_info = {}
        with open(summary_path, "r") as f:
            lines = f.readlines()

        for line in lines:
            line = line.strip()
            if line.startswith("Run ID:"):
                run_info["id"] = line.split(":", 1)[1].strip()
            elif line.startswith("Report Date:"):
                run_info["date"] = line.split(":", 1)[1].strip()
            elif line.startswith("Report A Status:"):
                run_info["status_a"] = parse_status_line(line, "DayViol")
            elif line.startswith("Report B Status:"):
                run_info["status_b"] = parse_status_line(line, "MonthViol")
            elif line.startswith("OverallStatus:"):
                run_info["overall"] = line.split(":", 1)[1].strip()
            elif line.startswith("ExceptionFlag:"):
                run_info["exception_flag"] = line.split(":", 1)[1].strip()

        # Fallback: compute overall if not in summary
        if "overall" not in run_info:
            stat_a = run_info.get("status_a", {}).get("status", "Pass")
            stat_b = run_info.get("status_b", {}).get("status", "Pass")
            if stat_a == "Violation" or stat_b == "Violation":
                run_info["overall"] = "Exception"
            else:
                run_info["overall"] = "Completed"

        runs.append(run_info)
    return runs


def parse_status_line(line, marker_name):
    """Parses a status line from summary.txt.

    Example: Report A Status: Pass (DayViol: N, Thresholds: {'Threshold1': 70, ...})
    """
    pattern = rf"Status:\s*(\w+)\s*\({marker_name}:\s*(\w+),\s*Thresholds:\s*(\{{.*\}})\)"
    match = re.search(pattern, line)
    if match:
        status = match.group(1)
        marker = match.group(2)
        try:
            thresholds = eval(match.group(3))
        except:
            thresholds = {}
        return {"status": status, "marker": marker, "thresholds": thresholds}
    return {"status": "Pending", "marker": "N", "thresholds": {}}


def set_slide_title(slide, text, size=28, color=PRIMARY):
    """Utility to set the title text and format it."""
    box = slide.shapes.title
    box.text = text
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color


def create_findings_deck(runs, output_file):
    """Generates the consolidated findings PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === 1. Title Slide ===
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ISAAI — Consolidated Validation Findings"
    slide.placeholders[1].text = (
        f"Automated analysis of {len(runs)} daily report validation runs\n"
        f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    )

    # === 2. Key Metrics Slide ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Validation Run Metrics")

    total = len(runs)
    completed = sum(1 for r in runs if r.get("overall") == "Completed")
    exceptions = sum(1 for r in runs if r.get("overall") == "Exception")

    # Count individual report statuses
    pass_a = sum(1 for r in runs if r.get("status_a", {}).get("status") == "Pass")
    viol_a = sum(1 for r in runs if r.get("status_a", {}).get("status") == "Violation")
    pass_b = sum(1 for r in runs if r.get("status_b", {}).get("status") == "Pass")
    viol_b = sum(1 for r in runs if r.get("status_b", {}).get("status") == "Violation")

    tf = slide.placeholders[1].text_frame
    tf.clear()

    metrics = [
        (f"Total Runs Analyzed: {total}", PRIMARY, True),
        (f"", None, False),
        (f"• Completed (no violations): {completed} ({completed/total*100:.1f}%)", PASS_GREEN, False),
        (f"• Exception (escalation required): {exceptions} ({exceptions/total*100:.1f}%)", ALERT_RED, True),
        (f"", None, False),
        (f"Report A (Daily — DayViol):", SECONDARY, True),
        (f"  • Pass: {pass_a}  |  Violation: {viol_a}", MUTED, False),
        (f"Report B (Monthly — MonthViol):", SECONDARY, True),
        (f"  • Pass: {pass_b}  |  Violation: {viol_b}", MUTED, False),
    ]

    for i, (text, color, bold) in enumerate(metrics):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        if color:
            p.font.color.rgb = color
        p.font.size = Pt(18)
        p.font.bold = bold

    # === 3. Run Ledger Table ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    title_box.text_frame.text = f"Validation Run Ledger ({min(15, len(runs))} most recent)"
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = PRIMARY

    display_runs = runs[-15:]  # Show last 15
    rows = len(display_runs) + 1
    cols = 6
    left, top, width, height = Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.0)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(0.8)   # Run ID
    table.columns[1].width = Inches(1.4)   # Date
    table.columns[2].width = Inches(1.8)   # Report A (DayViol)
    table.columns[3].width = Inches(1.8)   # Report B (MonthViol)
    table.columns[4].width = Inches(1.8)   # Overall Status
    table.columns[5].width = Inches(1.8)   # Exception Flag

    headers = ["Run", "Report Date", "Report A\n(DayViol)", "Report B\n(MonthViol)",
               "Overall Status", "Exception\nFlag"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, run in enumerate(display_runs):
        r_idx = row_idx + 1
        data = [
            run.get("id", "?"),
            run.get("date", "?"),
            run.get("status_a", {}).get("status", "Pending"),
            run.get("status_b", {}).get("status", "Pending"),
            run.get("overall", "Processing"),
            run.get("exception_flag", "No"),
        ]
        for c_idx, text in enumerate(data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

            if text == "Violation" or text == "Exception" or text == "Yes":
                p.font.color.rgb = ALERT_RED
                p.font.bold = True
            elif text == "Pass" or text == "Completed" or text == "No":
                p.font.color.rgb = PASS_GREEN

    # === 4. Exception Detail Slides ===
    exc_runs = [r for r in runs if r.get("overall") == "Exception"]
    if exc_runs:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_title(slide, f"Exception Deep-Dive ({len(exc_runs)} runs)", color=ALERT_RED)

        tf = slide.placeholders[1].text_frame
        tf.clear()

        for i, run in enumerate(exc_runs[:8]):  # Limit to 8 per slide
            run_id = run.get("id", "?")
            date = run.get("date", "?")
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"Run {run_id} ({date}):"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = ALERT_RED

            # Report A details
            sa = run.get("status_a", {})
            pa = tf.add_paragraph()
            pa.text = f"  Report A (DayViol={sa.get('marker', '?')}): {sa.get('status', '?')}"
            pa.font.size = Pt(14)
            pa.font.color.rgb = ALERT_RED if sa.get("status") == "Violation" else PASS_GREEN

            if sa.get("status") == "Violation":
                for t_name, t_val in sa.get("thresholds", {}).items():
                    if t_val < THRESHOLD_LIMIT:
                        pt = tf.add_paragraph()
                        pt.text = f"    ⚠ {t_name}: {t_val} (limit: {THRESHOLD_LIMIT})"
                        pt.font.size = Pt(12)
                        pt.font.color.rgb = ALERT_RED

            # Report B details
            sb = run.get("status_b", {})
            pb = tf.add_paragraph()
            pb.text = f"  Report B (MonthViol={sb.get('marker', '?')}): {sb.get('status', '?')}"
            pb.font.size = Pt(14)
            pb.font.color.rgb = ALERT_RED if sb.get("status") == "Violation" else PASS_GREEN

            if sb.get("status") == "Violation":
                for t_name, t_val in sb.get("thresholds", {}).items():
                    if t_val < THRESHOLD_LIMIT:
                        pt = tf.add_paragraph()
                        pt.text = f"    ⚠ {t_name}: {t_val} (limit: {THRESHOLD_LIMIT})"
                        pt.font.size = Pt(12)
                        pt.font.color.rgb = ALERT_RED

    # === 5. Summary Slide ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Conclusion")
    tf = slide.placeholders[1].text_frame
    tf.clear()

    conclusions = [
        f"Total validation runs analyzed: {total}",
        f"Successful completions: {completed} ({completed/total*100:.1f}%)",
        f"Escalations triggered: {exceptions} ({exceptions/total*100:.1f}%)",
        "",
        "The automated validation engine correctly identifies:",
        "  • DayViol / MonthViol markers in XML reports",
        "  • Threshold breaches below the configured limit (50)",
        "  • Exception cases requiring supervisor review",
        "",
        "Evidence files (XLSX) are generated for each report and stored in SharePoint.",
    ]

    for i, text in enumerate(conclusions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED

    # Save
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_file)
    print(f"Findings presentation generated: {output_file.relative_to(ROOT)}")


def main():
    PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)

    print("Scanning validation run results...")
    runs = parse_runs()

    if not runs:
        print("No validation runs found. Please generate some using scripts/simulate_runs.py first.")
        sys.exit(1)

    print(f"Found {len(runs)} validation runs.")

    output_path = PRESENTATIONS_DIR / "ISAAI-Consolidated-Findings.pptx"
    create_findings_deck(runs, output_path)

    print("Findings presentation compilation complete.")


if __name__ == "__main__":
    main()
