#!/usr/bin/env python3
"""ISAAI Automated PowerPoint presentation generator.

Reads the test runs from the local tests/runs/ directory and generates
consolidated slide decks and detailed reports in the presentations/Automated Presentations/
folder. Uses the workspace color theme and styling conventions.
"""

import sys
import re
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Setup paths relative to the script location
ROOT = Path(__file__).resolve().parents[1]
TESTS_DIR = ROOT / "tests"
RUNS_DIR = TESTS_DIR / "runs"
PRESENTATIONS_DIR = ROOT / "presentations" / "Automated Presentations"

THRESHOLD_LIMIT = 50

# Style guide colors matching the project charter
PRIMARY = RGBColor(0x00, 0x33, 0x5B)      # Slate Blue
SECONDARY = RGBColor(0x00, 0x50, 0x8F)    # Active Blue
MUTED = RGBColor(0x66, 0x66, 0x66)        # Neutral Gray
ALERT_RED = RGBColor(0xD9, 0x38, 0x3A)    # Warning Red
PASS_GREEN = RGBColor(0x2E, 0x7D, 0x32)   # Pass Green

def parse_runs():
    """Scans tests/runs/ directories, parsing run summaries."""
    runs = []
    if not RUNS_DIR.exists():
        print(f"Error: {RUNS_DIR} does not exist. Run scripts/simulate_runs.py first.")
        sys.exit(1)
        
    for run_path in sorted(RUNS_DIR.glob("run_*")):
        summary_path = run_path / "summary.txt"
        if not summary_path.exists():
            continue
            
        run_info = {}
        with open(summary_path, "r") as f:
            lines = f.readlines()
            
        for line in lines:
            line = line.strip()
            if line.startswith("Run ID:"):
                run_info["id"] = line.split(":", 1)[1].strip()
            elif line.startswith("Report Date:"):
                run_info["date"] = line.split(":", 1)[1].strip()
            elif line.startswith("Report A Status:"):
                run_info["status_a"] = parse_status_line(line)
            elif line.startswith("Report B Status:"):
                run_info["status_b"] = parse_status_line(line)
                
        # Overall outcome: if either has Exception, it's Exception.
        # If either has NoEscalation (but no Exception), it's NoEscalation.
        # Otherwise Happy.
        stat_a = run_info.get("status_a", {}).get("status")
        stat_b = run_info.get("status_b", {}).get("status")
        
        if stat_a == "Exception" or stat_b == "Exception":
            run_info["overall"] = "Exception"
        elif stat_a == "NoEscalation" or stat_b == "NoEscalation":
            run_info["overall"] = "NoEscalation"
        else:
            run_info["overall"] = "Happy"
            
        runs.append(run_info)
    return runs

def parse_status_line(line):
    """Parses a status line from summary.txt.
    
    Example: Report A Status: Exception (Marker: V, Thresholds: {'Threshold1': 70, 'Threshold2': 18, ...})
    """
    match = re.search(r"Status:\s*(\w+)\s*\(Marker:\s*(\w+),\s*Thresholds:\s*(\{.*\})\)", line)
    if match:
        status = match.group(1)
        marker = match.group(2)
        try:
            thresholds = eval(match.group(3))
        except:
            thresholds = {}
        return {"status": status, "marker": marker, "thresholds": thresholds}
    return {"status": "Unknown", "marker": "N", "thresholds": {}}

def set_slide_title(slide, text, size=28, color=PRIMARY):
    """Utility to set the title text and format it."""
    box = slide.shapes.title
    box.text = text
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color

def create_summary_deck(runs, output_file):
    """Generates the main PowerPoint presentation representing the test outcomes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ISAAI Run Report & Audit Deck"
    slide.placeholders[1].text = f"Automated summary of {len(runs)} simulated daily validation runs\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # 2. Key Metrics Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Validation Run Metrics")
    
    total = len(runs)
    happy = sum(1 for r in runs if r["overall"] == "Happy")
    no_esc = sum(1 for r in runs if r["overall"] == "NoEscalation")
    exceptions = sum(1 for r in runs if r["overall"] == "Exception")
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = f"Total Inbound Runs Analyzed: {total}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = f"• Happy Path (No violations detected): {happy} ({happy/total*100:.1f}%)"
    p2.font.size = Pt(18)
    p2.font.color.rgb = PASS_GREEN
    
    p3 = tf.add_paragraph()
    p3.text = f"• Warning (Violation marker present, thresholds pass): {no_esc} ({no_esc/total*100:.1f}%)"
    p3.font.size = Pt(18)
    p3.font.color.rgb = SECONDARY
    
    p4 = tf.add_paragraph()
    p4.text = f"• Escalate / Exception (Violation marker present, threshold breached): {exceptions} ({exceptions/total*100:.1f}%)"
    p4.font.size = Pt(18)
    p4.font.color.rgb = ALERT_RED
    p4.font.bold = True

    # 3. Overview Table Slide (Maximum 10 rows for clean fit, or multiple tables)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Title manually for blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    title_box.text_frame.text = "Recent Run Ledger (Last 10 Runs)"
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    # Add Table
    rows = min(11, len(runs) + 1)
    cols = 5
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9.0)
    height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Column Widths
    table.columns[0].width = Inches(1.2) # Run ID
    table.columns[1].width = Inches(1.8) # Date
    table.columns[2].width = Inches(2.0) # Report A
    table.columns[3].width = Inches(2.0) # Report B
    table.columns[4].width = Inches(2.0) # Overall Outcome
    
    # Set headers
    headers = ["Run ID", "Report Date", "Report Type A", "Report Type B", "Overall Status"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
    # Populate table with last 10 runs
    for row_idx, run in enumerate(runs[-10:]):
        r_idx = row_idx + 1
        data = [
            f"Run {run['id']}",
            run['date'],
            run['status_a']['status'],
            run['status_b']['status'],
            run['overall']
        ]
        for c_idx, text in enumerate(data):
            cell = table.cell(r_idx, c_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            
            # Format text color depending on status
            if text in ["Exception", "Breached"]:
                p.font.color.rgb = ALERT_RED
                p.font.bold = True
            elif text == "Happy" or text == "Pass":
                p.font.color.rgb = PASS_GREEN
            elif text == "NoEscalation":
                p.font.color.rgb = SECONDARY

    # 4. Exception Details Slides
    esc_runs = [r for r in runs if r["overall"] == "Exception"]
    for run in esc_runs[:5]: # limit to first 5 details slides to avoid bloat
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_title(slide, f"Exception Deep-Dive: Run {run['id']} ({run['date']})", color=ALERT_RED)
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Escalation triggered due to threshold breaches:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        # Report A details
        pa = tf.add_paragraph()
        pa.text = f"Report A status: {run['status_a']['status']}"
        pa.font.size = Pt(16)
        if run['status_a']['status'] == "Exception":
            pa.font.color.rgb = ALERT_RED
            for t_name, t_val in run['status_a']['thresholds'].items():
                p_t = tf.add_paragraph()
                status_str = "BREACHED" if t_val < THRESHOLD_LIMIT else "Pass"
                p_t.text = f"   - {t_name}: {t_val} (Limit: {THRESHOLD_LIMIT}) → {status_str}"
                p_t.font.size = Pt(14)
                p_t.font.color.rgb = ALERT_RED if t_val < THRESHOLD_LIMIT else MUTED
        else:
            pa.font.color.rgb = PASS_GREEN

        # Report B details
        pb = tf.add_paragraph()
        pb.text = f"\nReport B status: {run['status_b']['status']}"
        pb.font.size = Pt(16)
        if run['status_b']['status'] == "Exception":
            pb.font.color.rgb = ALERT_RED
            for t_name, t_val in run['status_b']['thresholds'].items():
                p_t = tf.add_paragraph()
                status_str = "BREACHED" if t_val < THRESHOLD_LIMIT else "Pass"
                p_t.text = f"   - {t_name}: {t_val} (Limit: {THRESHOLD_LIMIT}) → {status_str}"
                p_t.font.size = Pt(14)
                p_t.font.color.rgb = ALERT_RED if t_val < THRESHOLD_LIMIT else MUTED
        else:
            pb.font.color.rgb = PASS_GREEN

    # Save presentation
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_file)
    print(f"PowerPoint summary deck generated: {output_file.relative_to(ROOT)}")

def main():
    PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)
    
    print("Scanning validation run results...")
    runs = parse_runs()
    
    if not runs:
        print("No validation runs found. Please generate some using scripts/simulate_runs.py first.")
        sys.exit(1)
        
    print(f"Found {len(runs)} validation runs.")
    
    output_path = PRESENTATIONS_DIR / "ISAAI-Runs-Summary-Deck.pptx"
    create_summary_deck(runs, output_path)
    
    print("PowerPoint compilation complete.")

if __name__ == "__main__":
    main()
