#!/usr/bin/env python3
"""Build ISAAI boardroom PowerPoint (bookend narrative)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentations" / "ISAAI-Boardroom-Manual-to-Automated.pptx"

OPENING_QUESTION = (
    "Can we eliminate 30–40 minutes of daily manual report work—"
    "and still strengthen auditability and control?"
)

CLOSING_ANSWER_BULLETS = [
    "Yes—with scope discipline: automate ingest → validate → archive; humans only for exception governance.",
    "ISA: enterprise ArchiMate alignment, invoice PDF/OCR, chart automation.",
    "A+I: operable BPMN, XML/GOAL validation, OpenAPI integration mocks.",
    "Open: production LLM gateway, DMS access, sandbox account (see compliance docs).",
]

PRIMARY = RGBColor(0x00, 0x33, 0x5B)
SECONDARY = RGBColor(0x00, 0x50, 0x8F)
MUTED = RGBColor(0x66, 0x66, 0x66)


def set_title(slide, text, size=32, color=PRIMARY, center=False):
    box = slide.shapes.title
    box.text = text
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    if center:
        p.alignment = PP_ALIGN.CENTER


def add_bullets(slide, items, top=1.6, left=0.8, width=8.5, height=4.5, size=18):
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED


def add_image_slide(prs, title, image_path, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    path = ROOT / image_path
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.5), Inches(1.2), width=Inches(9))
    if subtitle:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tx.text_frame.text = subtitle
        tx.text_frame.paragraphs[0].font.size = Pt(12)
        tx.text_frame.paragraphs[0].font.color.rgb = MUTED


def question_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = OPENING_QUESTION
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    note = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(8.8), Inches(0.8))
    note.text_frame.text = "Pause. Do not answer yet."
    note.text_frame.paragraphs[0].font.size = Pt(14)
    note.text_frame.paragraphs[0].font.color.rgb = SECONDARY
    note.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def answer_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, "The answer", size=30)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = OPENING_QUESTION
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = SECONDARY
    for bullet in CLOSING_ANSWER_BULLETS:
        bp = tf.add_paragraph()
        bp.text = bullet
        bp.level = 0
        bp.font.size = Pt(17)
        bp.font.color.rgb = MUTED


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(s, "ISAAI", size=40)
    s.placeholders[1].text = "Dual-module automation program\nISA · Architecture & Integration"

    # 2 Hook
    question_slide(prs)

    # 3–6 Content
    slides_text = [
        ("Why this matters", [
            "30–40 minutes daily manual effort on structured XML reports",
            "Error risk: import, wrong field/version, inconsistent outputs",
            "Archive gap when steps are skipped",
        ]),
        ("Manual labor today (IST)", [
            "Phase 1: Mail → ZIP → XML extraction",
            "Phase 2: Spreadsheet import and threshold checks",
            "Phase 3: Evidence + distribution mail",
            "Phase 4: Document repository sign-off",
        ]),
        ("Risk & control gap", [
            "Media break: collaboration storage vs. document repository",
            "Upload → sign-off not system-enforced",
            "Naming convention friction = highest automation leverage",
        ]),
        ("Automation vision (SOLL)", [
            "Automated ingest, rule engine, exception queue",
            "Supervisor only on exceptions (human-in-the-loop)",
            "Standardized audit trail per event",
        ]),
    ]
    for title, bullets in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_title(slide, title)
        add_bullets(slide, bullets)

    # 7–8 BPMN
    add_image_slide(
        prs,
        "BPMN GOAL — automated main flow",
        "docs/Documentation/Process Modeling/GOAL Process/daily_financial_report_validation_lanes_OPERATIVE GOAL.png",
        "Ingest → validate → archive; automation lanes",
    )
    add_image_slide(
        prs,
        "BPMN GOAL — exception governance",
        "docs/Documentation/Architecture/GOAL Architecture/GOAL - Governance, Data and Motivation.png",
        "Exception path and audit-trail goal",
    )

    # 9–11 Modules
    module_slides = [
        ("How ISA shapes this project", [
            "Enterprise alignment via ArchiMate",
            "Exchange invoice: PDF/OCR → JSON/CSV → charts",
            "Delivers presentation-ready visual evidence",
        ]),
        ("How A+I shapes this project", [
            "Operable BPMN and XML/GOAL validation",
            "Integration architecture and white-box APIs",
            "SharePoint/DMS mocks for archival simulation",
        ]),
        ("ISA vs A+I — influence matrix", [
            "ISA: structure, alignment, invoice track",
            "A+I: process operability, financial report track",
            "Shared repo; separate branches and acceptance criteria",
        ]),
    ]
    for title, bullets in module_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_title(slide, title)
        add_bullets(slide, bullets)

    add_image_slide(
        prs,
        "Architecture IST → GOAL",
        "docs/Documentation/Architecture/GOAL Architecture/GOAL - Automated Main Flow.png",
    )

    for title, bullets in [
        ("Compliance & audit trail", [
            "Event model: who, when, action, source hash, approval state",
            "Pseudonymize before internal LLM (project default)",
            "OpenAPI v3 mocks document intended interfaces",
        ]),
        ("Roadmap to submission", [
            "M2 ISA module · M3 A+I module · M4 integration",
            "M5: models, docs, boardroom deck",
            "Target: 16 June 2026",
        ]),
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_title(slide, title)
        add_bullets(slide, bullets)

    answer_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
